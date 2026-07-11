from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()

# Standard colors
TITLE_COLOR = RGBColor(0, 70, 139) # Nature Blue
ACCENT_COLOR = RGBColor(237, 0, 0) # Nature Red
BODY_COLOR = RGBColor(60, 60, 60)

def set_font(font, name='Helvetica', size=None, color=None, bold=False):
    font.name = name
    if size:
        font.size = Pt(size)
    if color:
        font.color.rgb = color
    font.bold = bold

def add_title_slide(title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    set_font(title.text_frame.paragraphs[0].font, size=40, color=TITLE_COLOR, bold=True)
    
    subtitle.text = subtitle_text
    for p in subtitle.text_frame.paragraphs:
        set_font(p.font, size=24, color=BODY_COLOR)

def add_content_slide(title_text, content_list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = title_text
    set_font(title.text_frame.paragraphs[0].font, size=32, color=TITLE_COLOR, bold=True)
    
    tf = body.text_frame
    tf.clear()
    for text in content_list:
        p = tf.add_paragraph()
        p.text = text
        set_font(p.font, size=18, color=BODY_COLOR)

def add_image_slide(title_text, img_path, caption_list=None):
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    set_font(title.text_frame.paragraphs[0].font, size=32, color=TITLE_COLOR, bold=True)
    
    # Add image
    margin_top = Inches(1.5)
    margin_left = Inches(1.0)
    width = Inches(8)
    slide.shapes.add_picture(img_path, margin_left, margin_top, width=width)
    
    # Add caption
    if caption_list:
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(8), Inches(1))
        tf = txBox.text_frame
        for line in caption_list:
            p = tf.add_paragraph()
            p.text = line
            set_font(p.font, size=16, color=BODY_COLOR)

# Generate 30 slides content
slides_data = [
    # --- Intro ---
    ("title", "MT-LNN: A Brain-Inspired Liquid Neural Network\n液态神经网络重塑大模型架构", "Breaking the Memory Wall (打破内存墙)\nEverestAn | 2026", None),
    ("content", "1. The AI Memory Crisis (AI 行业的内存危机)", ["- LLMs compete on Context Length (200K+ tokens).", "- 但是，支撑超长文本的计算成本呈指数级爆炸。", "- OOM (Out Of Memory) limits local deployments.", "- Transformer's bottleneck is physical memory, not logic."]),
    ("content", "2. The Root Cause: KV Cache (元凶：KV缓存)", ["- Transformer acts like a forced video recorder.", "- 预测下一个词时，必须把过去每一个字的特征放在显存里。", "- Memory complexity scales at O(N). Compute scales at O(N^2).", "- 10万字的上下文会让一张昂贵的 A100 直接塞满。"]),
    ("content", "3. A Dead End for Enterprises (算力堆叠是死胡同)", ["- Cloud costs for millions of tokens are prohibitive.", "- 面对高并发请求，云服务商只能靠堆叠万亿算力。", "- Latency drops exponentially as context grows.", "- 端侧 (Edge AI) 无法运行大模型，因为手机 RAM 只有十几 GB。"]),
    ("content", "4. Biological Inspiration (生物学启示：人脑机制)", ["- The human brain does not remember every pixel.", "- 人脑绝不可能记住十年来每个画面的精确参数。", "- We rely on Working Memory and Selective Forgetting.", "- 我们依靠“工作记忆”和“选择性遗忘”。"]),
    
    # --- MT-LNN Core ---
    ("content", "5. What is a Liquid Neural Network? (什么是液态脑)", ["- An architecture whose state constantly dynamically flows.", "- 一种状态随时间动态流动的结构，区别于静态权重。", "- Forget the noise, compress the needle.", "- 只提取核心概要，遗忘冗余废话，把千言万语压缩。"]),
    ("content", "6. Introducing MT-LNN (我们的方案：MT-LNN 架构登场)", ["- Memory-Token Linear Neural Network.", "- 记忆令牌线性神经网络。", "- Replaces standard Self-Attention with Linear Attention.", "- 抛弃标准自注意力机制，采用纯线性算子。"]),
    ("content", "7. The Core: Selective Copy (核心：选择性记忆机制)", ["- Uses dynamic gating to explicitly select what to retain.", "- 引入动态门控，明确选择在长文中保留什么，丢弃什么。", "- Resolves 'Lost in the middle' phenomenon.", "- 彻底解决大模型长文中间的“遗忘”现象。"]),
    ("content", "8. Deep Dive 1: Parallel Scan (深度解析：线性并行扫描)", ["- Allows for highly efficient hardware utilization.", "- 并行扫描使得线性处理速度比传统 RNN 快百倍。", "- Every token adds O(1) constant overhead.", "- 每一个新单词的计算量都是 O(1) 常数，永远不升。"]),
    ("content", "9. Deep Dive 2: Quantum Coupling (深度解析：量子启发耦合)", ["- State mixing inspired by quantum entanglement states.", "- 使用量子态灵感混合隐藏状态。", "- Multi-head parallel routing of memory states.", "- 多头并行路由，增强特征维度的提取能力。"]),
    ("content", "10. The Mathematics of Edge Support (边缘计算的数学保障)", ["- State h_t = f(h_t-1, x_t).", "- Only the fixed-size h vector is kept in memory.", "- 模型不需要保存过去词的 KV 矩阵，仅需保留固定大小向量。", "- 完美适配 8GB 以下内存的手机和 AR 设备。"]),
    
    # --- Charts & Benchmarks ---
    ("image", "11. Cost Explosion vs. MT-LNN (成本爆炸图解)", "notes/fig_cost_scaling.png", ["Red: Transformer O(N) memory / O(N^2) compute.", "Blue: MT-LNN O(1) memory / O(1) compute.", "红线：Transformer长文本深渊；蓝线：MT-LNN的类脑常数飞行。"]),
    ("content", "12. Benchmark 1: Needle in a Haystack (核心评测：大海捞针)", ["- Tested on 128K context retrieval.", "- 12万字极度长文提取测试。", "- MT-LNN retains 99%+ accuracy independent of depth.", "- 无视深度，精准命中（得益于 Selective Copy）。"]),
    ("content", "13. Benchmark 2: Selective Copying (核心评测：抗干扰测试)", ["- Retaining sequences separated by 1000+ noisy tokens.", "- 跨越千字噪音阻隔，精准复制指令。", "- Completely outperforms classic global attention.", "- 面对巨量噪音的稳健性碾压传统注意力。"]),
    ("image", "14. Performance Radar (雷达图：维度对比)", "notes/fig_benchmark_radar.png", ["MT-LNN matches SOTA in Long Context & blows them away in Deployment/Cost.", "通用知识因参数量小暂处劣势，但在核心部署性能上碾压大厂。"]),
    ("content", "15. The Honest Trade-off (诚实的权衡：放弃百科知识)", ["- We are not competing to be a Wikipedia engine.", "- 我们不和千亿美金的 GPT-4o 拼“通用常识百科库”。", "- We focus on reasoning logic and ultra-long structural memory.", "- 把所有参数压在逻辑推理与无限期超级记忆上。"]),
    ("content", "16. Competitor Analysis: Claude/GPT (对标巨头：垂直收割)", ["- Claude 3.5 requires API/Cloud, zero privacy.", "- 大厂模型无法本地化，企业数据缺乏隐私。", "- MT-LNN offers private, localized million-token RAG.", "- MT-LNN 提供完全本地化、隐私受控的长文知识库。"]),
    ("content", "17. Competitor Analysis: Mamba & RWKV (生态对抗：同赛道)", ["- Open-source linear models lack needle-retrieval fidelity.", "- Mamba 等方案难以保证极长文精确寻点。", "- MT-LNN's explicit state tracking solves linear deterioration.", "- 我们的显式量子门控阻止了线性衰减。"]),
    
    # --- Economics & Deploy ---
    ("content", "18. Economics: 90% Cost Reduction (商业优势重塑：降本 90%)", ["- Cloud hosting inference drops from millions to thousands.", "- 云端推理服务器架构大幅度降级，每月节约数万元。", "- Server costs no longer tie to context length.", "- 用户不管输入多长的文本，不再增加任何云端边际成本。"]),
    ("image", "19. High Concurrency Capacity (高并发承载力)", "notes/fig_concurrency.png", ["To handle 100 users querying 100k tokens at the same time:", "A100 crashes instantly. RTX 4090 handles it with MT-LNN's constant state."]),
    ("content", "20. Market 1: Edge Computing & IoT (下沉市场：IoT与端侧)", ["- Phones, AR Glasses, On-device assistants.", "- 手机私有大模型、AR眼镜全天候音视频分析核心。", "- Continuous streaming inputs won't drain the battery.", "- 绝不会因为 24 小时保持倾听而撑爆内存、发烫。"]),
    ("content", "21. Market 2: B2B Enterprise Deploy (企业痛点：私有化部署)", ["- Law firms, financial PDF audits, codebase analysis.", "- 律所看卷宗、投行看财报、科技公司分析百万行代码。", "- Requires large context, requires privacy, requires low hardware.", "- 需要长大模型，且数据不能出域。完美契合 MT-LNN。"]),
    ("image", "22. The Future Tech Landscape (技术生态排位占领)", "notes/fig_landscape.png", ["Top Right: The ultimate quadrant of long context AND edge deployability.", "右上角：唯一兼具极长文本掌控力与下沉端侧部署能力的架构。"]),
    
    # --- Roadmap ---
    ("content", "23. Strategic Roadmap Overview (三阶段路线图概览)", ["- Moving from Proof of Concept to General Artificial Intelligence.", "- 从本地超级知识库，走向下一代通用大模型。", "- Fast, extremely capital efficient progression.", "- 凭借架构极度省算力的优势，拥有最快的迭代速度。"]),
    ("content", "24. Roadmap Stage 1: The 1.5B Proof (阶段一：验证原语)", ["- Timeline: Now.", "- 现状：已用 15 美元，在 1.5B Qwen 规模上跑通极长大海捞针。", "- Goal: Secure initial VC funding and publish core paper.", "- 目标：通过极小成本的 Demo 验证原设。"]),
    ("content", "25. Roadmap Stage 2: The 7B King (阶段二：专业微调王者)", ["- Timeline: Next 3-6 months.", "- 未来半载：使用 4x A100 训练 7B 企业级原生模型。", "- Application: Sell vertical B2B solutions.", "- 应用：全面碾压企业级合同审查、金融报告等长文档 RAG。"]),
    ("content", "26. Roadmap Stage 3: AGI Challenger (阶段三：通用架构颠覆)", ["- Timeline: 1 to 3 Years.", "- 终局之战：筹建 512x H100 集群，训出 70B+ 千亿巨兽。", "- Full-scale replacement of standard Transformer architectures.", "- 在全范围（含通用知识）内，全面替换过时的 Transformer。"]),
    ("content", "27. Defensibility & IP Moats (护城河在哪里？)", ["- Core algorithmic code base open-sourced for standard adoption.", "- 算法基础开源，吸引广大开发者形成生态。", "- Specialized quantization and CUDA kernels remain in-house.", "- 高速并行 CUDA 算子与极限压缩量化闭源提供商业支撑。"]),
    ("content", "28. Current Progress & RAG UI Demo (目前应用：RAG 系统)", ["- Successfully built a complete RAG UI in PyTorch/Gradio.", "- 已跑通完全本地部署的百万字个人知识库 RAG Demo。", "- Deploys immediately via container to any platform.", "- 一键部署支持云平台、HuggingFace 与本地 Windows。"]),
    
    # --- Outro ---
    ("content", "29. Summary: Why Invest Now? (结语：为何现在下注？)", ["- Transformer has reached its physical efficiency limit.", "- Transformer 路线“大力出奇迹”已达物理硬件的极限瓶颈。", "- The memory wall requires a fundamental architecture shift.", "- 只有底层算法级的跨越，才能打破昂贵的“内存墙”。", "- We have the architecture, the code, and the benchmarks.", "- 我们握有类脑算法钥匙、能跑出来的实机代码、与碾压级降本数据。"]),
    ("content", "30. Appendix & Links (资源与相关链接)", ["- Demo Repo: github.com/everest-an/Awareness-O1", "- Core Arch Repo: github.com/everest-an/O1", "- Cloud Training Setup: check CLOUD_TRAINING_GUIDE.md", "\nThank you. Liquidize the AI future."])
]

# Generate
for slide in slides_data:
    if slide[0] == "title":
        add_title_slide(slide[1], slide[2])
    elif slide[0] == "content":
        add_content_slide(slide[1], slide[2])
    elif slide[0] == "image":
        add_image_slide(slide[1], slide[2], slide[3])

prs.save("Investor_Pitch_Deck_MT_LNN.pptx")
print("Presentation generated successfully at Investor_Pitch_Deck_MT_LNN.pptx")
